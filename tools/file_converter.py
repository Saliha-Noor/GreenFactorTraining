import os
from pptx import Presentation
from docx import Document


def convert_pptx(filepath):
    print(f"  Converting PPTX: {os.path.basename(filepath)}")
    prs = Presentation(filepath)
    text = ""
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text.strip() + "\n"
        if slide_text:
            text += f"\n--- Slide {slide_num} ---\n{slide_text}"
    return text


def convert_docx(filepath):
    print(f"  Converting DOCX: {os.path.basename(filepath)}")
    doc = Document(filepath)
    text = ""
    for para in doc.paragraphs:
        if para.text.strip():
            text += para.text.strip() + "\n"
    return text


def convert_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".pptx":
        return convert_pptx(filepath)
    elif ext == ".docx":
        return convert_docx(filepath)
    elif ext == ".pdf":
        return None
    else:
        print(f"  Skipping unsupported file: {os.path.basename(filepath)}")
        return None